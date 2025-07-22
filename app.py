from flask import Flask, render_template, request, jsonify, send_from_directory
import os
import time
from werkzeug.utils import secure_filename
import whisper
import google.generativeai as genai
from moviepy.editor import VideoFileClip
import json
import random
import logging
from gtts import gTTS
import re
import numpy as np
import math
from ratelimit import limits, sleep_and_retry
# PPT analysis imports
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import uuid

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Rate limiting configuration
CALLS = 60  # Number of calls allowed
RATE_LIMIT_PERIOD = 60  # Time period in seconds

@sleep_and_retry
@limits(calls=CALLS, period=RATE_LIMIT_PERIOD)
def rate_limited_api_call(func):
    """Decorator to rate limit API calls"""
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            if "quota exceeded" in str(e).lower():
                logger.error("API quota exceeded. Waiting before retry...")
                time.sleep(5)  # Wait 5 seconds before retry
                return func(*args, **kwargs)
            raise e
    return wrapper

# Configure Gemini API
GOOGLE_API_KEY = 'AIzaSyB5FFNNqzfBeupJOQx3MnXkyUUBr7tAXeI'

def configure_gemini_api():
    """Configure and test the Gemini API connection"""
    try:
        # Configure the API
        genai.configure(api_key=GOOGLE_API_KEY)
        
        # Get available models
        models = genai.list_models()
        if not models:
            logger.error("No models available")
            return False
            
        # Try each model in order of preference
        for model_name in MODEL_NAMES:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content("Test connection")
                if response:
                    logger.info(f"Successfully configured Gemini API with model: {model_name}")
                    return True
            except Exception as e:
                logger.warning(f"Failed to use model {model_name}: {str(e)}")
                continue
        
        # If no preferred model works, try any Gemini model
        for model in models:
            if "gemini" in model.name.lower():
                try:
                    model = genai.GenerativeModel(model.name)
                    response = model.generate_content("Test connection")
                    if response:
                        logger.info(f"Successfully configured Gemini API with model: {model.name}")
                        return True
                except Exception as e:
                    logger.warning(f"Failed to use model {model.name}: {str(e)}")
                    continue
        
        logger.error("No working models found")
        return False
            
    except Exception as e:
        logger.error(f"API configuration failed: {str(e)}")
        return False

# Initialize API configuration at startup
if not configure_gemini_api():
    logger.error("Failed to initialize Gemini API. Please check your API key and quota.")
    print("ERROR: Failed to initialize Gemini API. Please check your API key and quota.")
    print("Please make sure:")
    print("1. Your API key is valid and has not expired")
    print("2. You have enabled the Gemini API in your Google Cloud Console")
    print("3. You have sufficient quota available")
    print("4. Your API key has the necessary permissions")
    print("5. You have a stable internet connection")
    print("\nDetailed error information has been logged.")
else:
    logger.info("Gemini API configured successfully")
    print("Gemini API configured successfully")

app = Flask(__name__, static_folder='static')

# Configure upload folders
UPLOAD_FOLDER = 'uploads'
AUDIO_FOLDER = 'audio'
TRANSCRIPTION_FOLDER = 'transcriptions'
OUTPUT_FOLDER = 'outputs'
TTS_FOLDER = 'tts_output'
IMPROVEMENTS_FOLDER = 'improvements'
USER_DATA_FOLDER = 'user_data'
PPT_UPLOAD_FOLDER = 'ppt_uploads'  # New folder for PPT files

# Ensure necessary directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)
os.makedirs(TRANSCRIPTION_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(TTS_FOLDER, exist_ok=True)
os.makedirs(IMPROVEMENTS_FOLDER, exist_ok=True)
os.makedirs(USER_DATA_FOLDER, exist_ok=True)
os.makedirs(PPT_UPLOAD_FOLDER, exist_ok=True)  # Create PPT upload folder

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload size

# Store uploaded PPT files
UPLOADED_PPT_FILES = {}

# List of preferred model names
MODEL_NAMES = [
    "models/gemini-1.5-pro-latest",
    "models/gemini-1.5-pro",
    "models/gemini-1.5-flash-latest",
    "models/gemini-1.5-flash",
    "models/gemini-2.0-pro-exp",
    "models/gemini-2.0-flash",
    "models/gemini-1.0-pro-vision-latest",
    "models/gemini-pro-vision"
]

# Common filler words to detect
FILLER_WORDS = [
    "um", "uh", "er", "ah", "like", "you know", "sort of", "kind of", 
    "basically", "actually", "literally", "so", "well", "i mean", "right"
]

# Skill evaluation criteria
SKILL_CRITERIA = {
    "filler_words": {
        "beginner": {
            "max_count": 5,
            "target_count": 3,
            "passing_score": 70
        },
        "intermediate": {
            "max_count": 3,
            "target_count": 1,
            "passing_score": 80
        },
        "advanced": {
            "max_count": 1,
            "target_count": 0,
            "passing_score": 90
        }
    },
    "tone_modulation": {
        "beginner": {
            "min_variation": 0.2,
            "target_variation": 0.3,
            "passing_score": 70
        },
        "intermediate": {
            "min_variation": 0.3,
            "target_variation": 0.4,
            "passing_score": 80
        },
        "advanced": {
            "min_variation": 0.4,
            "target_variation": 0.5,
            "passing_score": 90
        }
    },
    "posture": {
        "beginner": {
            "min_score": 0.6,
            "target_score": 0.7,
            "passing_score": 70
        },
        "intermediate": {
            "min_score": 0.7,
            "target_score": 0.8,
            "passing_score": 80
        },
        "advanced": {
            "min_score": 0.8,
            "target_score": 0.9,
            "passing_score": 90
        }
    },
    "eye_contact": {
        "beginner": {
            "min_duration": 0.5,  # Percentage of time maintaining eye contact
            "target_duration": 0.6,
            "passing_score": 70
        },
        "intermediate": {
            "min_duration": 0.6,
            "target_duration": 0.7,
            "passing_score": 80
        },
        "advanced": {
            "min_duration": 0.7,
            "target_duration": 0.8,
            "passing_score": 90
        }
    },
    "gestures": {
        "beginner": {
            "min_count": 2,
            "target_count": 3,
            "passing_score": 70
        },
        "intermediate": {
            "min_count": 3,
            "target_count": 5,
            "passing_score": 80
        },
        "advanced": {
            "min_count": 5,
            "target_count": 7,
            "passing_score": 90
        }
    },
    "pace": {
        "beginner": {
            "min_wpm": 120,
            "max_wpm": 150,
            "target_wpm": 130,
            "passing_score": 70
        },
        "intermediate": {
            "min_wpm": 130,
            "max_wpm": 160,
            "target_wpm": 145,
            "passing_score": 80
        },
        "advanced": {
            "min_wpm": 140,
            "max_wpm": 170,
            "target_wpm": 155,
            "passing_score": 90
        }
    },
    "clarity": {
        "beginner": {
            "min_score": 0.6,
            "target_score": 0.7,
            "passing_score": 70
        },
        "intermediate": {
            "min_score": 0.7,
            "target_score": 0.8,
            "passing_score": 80
        },
        "advanced": {
            "min_score": 0.8,
            "target_score": 0.9,
            "passing_score": 90
        }
    }
}

# Task prompts for each skill and level
TASK_PROMPTS = {
    "filler_words": {
        "beginner": "Record a 30-second introduction about yourself. Focus on eliminating filler words like 'um', 'like', and 'you know'.",
        "intermediate": "Deliver a 1-minute impromptu speech about your favorite movie. Challenge yourself to use minimal filler words while maintaining a natural flow.",
        "advanced": "Prepare a 2-minute formal presentation about a current event. Aim for zero filler words while maintaining professional composure."
    },
    "tone_modulation": {
        "beginner": "Read a short paragraph aloud with appropriate emotion. Practice varying your tone to match the content.",
        "intermediate": "Deliver a 1-minute speech about something you're passionate about, consciously varying your tone to emphasize key points.",
        "advanced": "Record a 2-minute persuasive speech, using tone modulation to convey different emotions and maintain audience engagement."
    },
    "posture": {
        "beginner": "Record yourself standing with proper posture for 30 seconds while introducing yourself. Focus on keeping your back straight and shoulders relaxed.",
        "intermediate": "Deliver a 1-minute speech while maintaining confident posture. Pay attention to your body alignment throughout.",
        "advanced": "Present a 2-minute speech while maintaining professional posture, incorporating natural movement and gestures."
    },
    "eye_contact": {
        "beginner": "Record a 30-second introduction while making eye contact with the camera. Try not to look away or down at notes.",
        "intermediate": "Deliver a 1-minute speech while maintaining consistent eye contact with the camera, practicing the 3-second rule.",
        "advanced": "Present a 2-minute speech while effectively managing eye contact, simulating looking at different audience members."
    },
    "gestures": {
        "beginner": "Record a 30-second speech while incorporating 2-3 basic hand gestures to emphasize key points.",
        "intermediate": "Deliver a 1-minute speech using purposeful gestures that complement your message.",
        "advanced": "Present a 2-minute speech with natural, varied gestures that enhance your message and demonstrate confidence."
    },
    "pace": {
        "beginner": "Record yourself reading a short paragraph at a steady, comfortable pace. Aim for about 120-150 words per minute.",
        "intermediate": "Deliver a 1-minute speech while consciously varying your pace to emphasize important points.",
        "advanced": "Present a 2-minute speech with dynamic pacing, using speed variations strategically for impact."
    },
    "clarity": {
        "beginner": "Record a 30-second speech about your favorite food, focusing on clear pronunciation and enunciation.",
        "intermediate": "Deliver a 1-minute speech about a complex topic, ensuring your words are clear and easily understood.",
        "advanced": "Present a 2-minute speech with specialized terminology while maintaining clarity and audience understanding."
    }
}

# PPT Analysis Functions
def extract_text_from_file(file_path):
    """Extract text from different file types"""
    file_extension = os.path.splitext(file_path)[1].lower()
    print(f"Processing file with extension: {file_extension}")
    
    if file_extension == '.pdf':
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    elif file_extension == '.docx':
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    elif file_extension == '.pptx':
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file_extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    else:
        raise ValueError(f"Unsupported file format: {file_extension}. Supported formats are .pdf, .docx, .pptx, and .txt.")

def extract_text_from_slide(file_path, slide_number):
    """Extract text from a specific slide in PPTX"""
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension != '.pptx':
        raise ValueError("Only PPTX files are supported for slide-specific analysis")
    
    prs = Presentation(file_path)
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(f"Slide number must be between 1 and {len(prs.slides)}")
    
    slide = prs.slides[slide_number - 1]
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
    return text, len(prs.slides)

def generate_speaking_script(text, speaker_name):
    """Generate a speaking script with real-time tips"""
    prompt = f"""
    Create a speaking script for {speaker_name} based on this presentation content. Include:
    1. A natural introduction
    2. Speaking tips marked with [TIP: your tip here]
    3. Emphasis points marked with [EMPHASIS: point to emphasize]
    4. Pause points marked with [PAUSE]
    5. Volume/pace adjustments marked with [LOUDER] or [SLOWER]
    6. A confident conclusion

    Content: {text}

    Format the response as:
    - Introduction: [script]
    - Main Content: [script with tips]
    - Conclusion: [script]
    """
    
    try:
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(prompt)
        if not response.text:
            return "Could not generate speaking script. Please try again."
        return response.text
    except Exception as e:
        print("Gemini API error:", str(e))
        return f"Error: Could not generate script. {str(e)}"

def generate_speaking_tips(text, speaker_name):
    """Generate speaking tips for a specific slide"""
    prompt = f"""
    Create speaking tips for {speaker_name} for this slide content. Include:
    1. How to introduce this slide
    2. Key points to emphasize
    3. When to pause
    4. When to speak louder/softer
    5. How to engage the audience
    6. How to transition to the next slide

    Slide Content: {text}

    Format the response as:
    - Introduction: [how to start]
    - Speaking Tips: [list of tips with [PAUSE], [LOUDER], [SOFTER] markers]
    - Key Points: [what to emphasize]
    - Audience Engagement: [how to engage]
    - Transition: [how to end and move to next slide]
    """
    
    try:
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(prompt)
        if not response.text:
            return "Could not generate speaking tips. Please try again."
        return response.text
    except Exception as e:
        print("Gemini API error:", str(e))
        return f"Error: Could not generate tips. {str(e)}"

# Existing functions (keeping all the original functions)
def get_available_models():
    """Get list of available models"""
    try:
        models = genai.list_models()
        return [model.name for model in models]
    except Exception as e:
        logger.error(f"Error listing models: {str(e)}")
        return []

def find_working_model():
    """Find a usable Gemini model"""
    try:
        # Get list of available models
        models = genai.list_models()
        available_models = [model.name for model in models]
        logger.info(f"Available models: {available_models}")

        # Try each model in order of preference
        for model_name in MODEL_NAMES:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content("Test connection")
                if response:
                    logger.info(f"Found working model: {model_name}")
                    return model_name
            except Exception as e:
                logger.warning(f"Failed to use model {model_name}: {str(e)}")
                continue

        # If no preferred model is found, try any Gemini model
        for model in available_models:
            if "gemini" in model.lower():
                try:
                    model = genai.GenerativeModel(model)
                    response = model.generate_content("Test connection")
                    if response:
                        logger.info(f"Using available Gemini model: {model}")
                        return model
                except Exception as e:
                    logger.warning(f"Failed to use model {model}: {str(e)}")
                    continue

        logger.warning("No working models found, defaulting to models/gemini-1.5-pro-latest")
        return "models/gemini-1.5-pro-latest"
    except Exception as e:
        logger.error(f"Error finding working model: {str(e)}")
        return "models/gemini-1.5-pro-latest"  # Default fallback

def get_model():
    """Get a configured Gemini model with retry logic"""
    max_retries = 3
    retry_delay = 5  # seconds
    
    for attempt in range(max_retries):
        try:
            model_name = find_working_model()
            model = genai.GenerativeModel(model_name)
            # Test the model
            response = model.generate_content("Test connection")
            if response:
                return model
        except Exception as e:
            logger.warning(f"Attempt {attempt + 1}/{max_retries} failed: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                continue
            raise Exception(f"Failed to initialize model after {max_retries} attempts: {str(e)}")
    
    raise Exception("Failed to initialize model")

def generate_confidence_scores():
    """Generate random confidence scores"""
    return {
        'transcription': random.uniform(0.85, 0.95),
        'content_analysis': random.uniform(0.80, 0.90),
        'body_language': {
            'facial_expressions': random.uniform(0.75, 0.85),
            'eye_contact': random.uniform(0.70, 0.80),
            'posture': random.uniform(0.75, 0.85),
            'gestures': random.uniform(0.70, 0.80),
            'overall': random.uniform(0.75, 0.85)
        },
        'speech_analysis': {
            'pace': random.randint(140, 180),  # Words per minute
            'vocal_variety': random.uniform(0.70, 0.90),
            'clarity': random.uniform(0.75, 0.90),
            'filler_words_count': random.randint(5, 15),
            'pauses_score': random.uniform(0.65, 0.85),
            'emphasis_score': random.uniform(0.70, 0.90),
            'pronunciation_score': random.uniform(0.75, 0.95),
            'overall': random.uniform(0.75, 0.90),
            'duration_seconds': random.randint(120, 300),
            'waveform_data': generate_random_waveform(100),
            'filler_words': []
        }
    }

def generate_random_waveform(length):
    """Generate random waveform data for visualization"""
    # Create a sine wave with some noise
    x = np.linspace(0, 4 * np.pi, length)
    base_wave = np.sin(x) * 0.5 + 0.5
    
    # Add some random noise
    noise = np.random.random(length) * 0.3
    
    # Combine and ensure values are between 0 and 1
    waveform = np.clip(base_wave + noise, 0, 1)
    
    # Add some emphasis points for visualization
    emphasis_points = []
    for i in range(3, 8):
        position = random.randint(0, length - 1)
        emphasis_points.append({
            "position": position,
            "type": "emphasis" if i % 2 == 0 else "pause"
        })
    
    return waveform.tolist()

def text_to_speech(text, filename):
    """Convert text to speech and save as MP3"""
    try:
        tts = gTTS(text=text, lang='en', slow=False)
        file_path = os.path.join(TTS_FOLDER, filename)
        tts.save(file_path)
        return file_path
    except Exception as e:
        logger.error(f"Error generating text-to-speech: {str(e)}")
        return None

def extract_key_points(analysis_text):
    """Extract key points from analysis text"""
    try:
        # Use Gemini to extract key points
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        prompt = f"""
        Extract 3-5 key points from the following analysis. 
        Format each point as a JSON object with 'point' and 'importance' (high, medium, low) properties:
        
        {analysis_text}
        
        Return ONLY a valid JSON array of objects.
        """
        
        response = model.generate_content(prompt)
        response_text = response.text
        
        # Extract JSON from response (handling potential formatting issues)
        json_start = response_text.find('[')
        json_end = response_text.rfind(']') + 1
        
        if json_start >= 0 and json_end > json_start:
            json_str = response_text[json_start:json_end]
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass
        
        # Fallback to simple extraction if JSON parsing fails
        lines = analysis_text.split('\n')
        key_points = []
        
        for line in lines:
            line = line.strip()
            if line and (line.startswith('-') or line.startswith('•') or (len(line) > 10 and '.' in line[:10])):
                key_points.append({
                    'point': line.lstrip('-•. '),
                    'importance': 'medium'
                })
                
                if len(key_points) >= 5:
                    break
                    
        return key_points[:5] if key_points else [{'point': 'Analysis completed', 'importance': 'medium'}]
        
    except Exception as e:
        logger.error(f"Error extracting key points: {str(e)}")
        return [{'point': 'Analysis completed', 'importance': 'medium'}]

# Function to identify improvements needed
def identify_improvements(transcription_text, body_analysis_text, speech_analysis_text):
    """Identify improvements needed in the presentation"""
    try:
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        prompt = f"""
        Based on the following transcription, body language analysis, and speech analysis, identify 5 specific improvements 
        the presenter could make. For each improvement, provide:
        1. A clear description of what needs improvement
        2. Why it's important
        3. Specific actionable advice on how to improve
        4. A priority level (high, medium, low)
        
        Format the response as a JSON array of objects with properties: 'area', 'importance', 'why', 'how_to_improve'
        
        Transcription:
        {transcription_text}
        
        Body Language Analysis:
        {body_analysis_text}
        
        Speech Analysis:
        {speech_analysis_text}
        
        Return ONLY a valid JSON array of objects.
        """
        
        response = model.generate_content(prompt)
        response_text = response.text
        
        # Extract JSON from response
        json_start = response_text.find('[')
        json_end = response_text.rfind(']') + 1
        
        if json_start >= 0 and json_end > json_start:
            json_str = response_text[json_start:json_end]
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass
        
        # Fallback if JSON parsing fails
        return [
            {
                'area': 'Content structure',
                'importance': 'high',
                'why': 'Well-structured content is easier to follow',
                'how_to_improve': 'Use clear introduction, body, and conclusion'
            },
            {
                'area': 'Delivery pace',
                'importance': 'medium',
                'why': 'Appropriate pace helps audience comprehension',
                'how_to_improve': 'Practice varying speech rate for emphasis'
            }
        ]
        
    except Exception as e:
        logger.error(f"Error identifying improvements: {str(e)}")
        return [
            {
                'area': 'General presentation',
                'importance': 'medium',
                'why': 'Effective presentation enhances message delivery',
                'how_to_improve': 'Review and practice presentation techniques'
            }
        ]

def analyze_speech(transcription_text, audio_path=None):
    """Analyze speech patterns and characteristics"""
    try:
        # Calculate speech metrics
        word_count = len(transcription_text.split())
        
        # Estimate duration (if audio_path is not available)
        duration_seconds = 0
        if audio_path and os.path.exists(audio_path):
            try:
                audio_clip = VideoFileClip(audio_path)
                duration_seconds = audio_clip.duration
                audio_clip.close()
            except Exception as e:
                logger.error(f"Error getting audio duration: {str(e)}")
                duration_seconds = word_count / 2.5  # Rough estimate: 150 wpm
        else:
            duration_seconds = word_count / 2.5  # Rough estimate: 150 wpm
        
        # Calculate words per minute
        words_per_minute = int(word_count / (duration_seconds / 60)) if duration_seconds > 0 else 150
        
        # Detect filler words
        filler_word_counts = {}
        for filler in FILLER_WORDS:
            # Use regex to find whole word matches
            pattern = r'\b' + re.escape(filler) + r'\b'
            matches = re.findall(pattern, transcription_text.lower())
            if matches:
                filler_word_counts[filler] = len(matches)
        
        # Convert to list format for frontend
        filler_words_list = [{"word": word, "count": count} for word, count in filler_word_counts.items()]
        filler_words_list.sort(key=lambda x: x["count"], reverse=True)
        
        # Get total filler word count
        total_filler_count = sum(filler_word_counts.values())
        
        # Use Gemini for more detailed speech analysis
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        prompt = f"""
        Analyze the following transcription for speech quality. Focus on:
        1. Vocal variety (changes in pitch, tone, volume)
        2. Clarity and pronunciation
        3. Effective use of pauses
        4. Emphasis on key points
        5. Overall speech effectiveness
        
        Provide specific examples from the text to support your analysis.
        
        Transcription:
        {transcription_text}
        """
        
        response = model.generate_content(prompt)
        speech_analysis_text = response.text
        
        # Generate waveform data for visualization
        waveform_length = 100  # Number of data points
        if audio_path and os.path.exists(audio_path):
            try:
                # In a real implementation, we would analyze the audio file
                # For this demo, we'll generate random data
                waveform_data = generate_random_waveform(waveform_length)
            except Exception as e:
                logger.error(f"Error generating waveform: {str(e)}")
                waveform_data = generate_random_waveform(waveform_length)
        else:
            waveform_data = generate_random_waveform(waveform_length)
        
        # Calculate scores based on analysis
        # In a real implementation, these would be derived from actual audio analysis
        vocal_variety_score = random.uniform(0.7, 0.9)
        clarity_score = random.uniform(0.75, 0.95)
        pauses_score = random.uniform(0.65, 0.85)
        emphasis_score = random.uniform(0.7, 0.9)
        pronunciation_score = random.uniform(0.75, 0.95)
        
        # Calculate overall score
        overall_score = (vocal_variety_score + clarity_score + pauses_score + emphasis_score + pronunciation_score) / 5
        
        # Return speech analysis results
        return {
            'text': speech_analysis_text,
            'metrics': {
                'pace': words_per_minute,
                'vocal_variety': vocal_variety_score,
                'clarity': clarity_score,
                'filler_words_count': total_filler_count,
                'filler_words': filler_words_list[:10],  # Limit to top 10
                'pauses_score': pauses_score,
                'emphasis_score': emphasis_score,
                'pronunciation_score': pronunciation_score,
                'overall': overall_score,
                'duration_seconds': duration_seconds,
                'waveform_data': waveform_data
            }
        }
        
    except Exception as e:
        logger.error(f"Error analyzing speech: {str(e)}")
        # Return default values if analysis fails
        return {
            'text': "Speech analysis could not be completed.",
            'metrics': {
                'pace': 150,
                'vocal_variety': 0.75,
                'clarity': 0.8,
                'filler_words_count': 0,
                'filler_words': [],
                'pauses_score': 0.7,
                'emphasis_score': 0.75,
                'pronunciation_score': 0.8,
                'overall': 0.75,
                'duration_seconds': 120,
                'waveform_data': generate_random_waveform(100)
            }
        }

def enhance_transcription(transcription_text):
    """Enhance transcription with improved vocabulary and grammar"""
    try:
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        prompt = f"""
        Enhance the following transcription by:
        1. Improving grammar and sentence structure
        2. Replacing basic vocabulary with more sophisticated alternatives where appropriate
        3. Maintaining the original meaning and speaker's intent
        4. Removing filler words and repetitions
        5. Organizing into clear paragraphs
        
        Original Transcription:
        {transcription_text}
        
        Return ONLY the enhanced transcription text without any explanations or meta-commentary.
        """
        
        response = model.generate_content(prompt)
        enhanced_text = response.text
        
        return enhanced_text
        
    except Exception as e:
        logger.error(f"Error enhancing transcription: {str(e)}")
        return transcription_text  # Return original if enhancement fails

def generate_all_suggestions_summary(improvements, body_analysis, speech_analysis):
    """Generate a comprehensive summary of all suggestions for text-to-speech"""
    try:
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        # Extract improvements data
        improvements_text = ""
        if improvements and isinstance(improvements, list):
            for i, imp in enumerate(improvements, 1):
                area = imp.get('area', 'Improvement area')
                importance = imp.get('importance', 'medium')
                why = imp.get('why', 'Important for effective communication')
                how = imp.get('how_to_improve', 'Practice and seek feedback')
                
                improvements_text += f"Suggestion {i}: {area}. Priority: {importance}. {why}. To improve: {how}.\n\n"
        
        prompt = f"""
        Create a concise, conversational summary of presentation improvement suggestions based on the following data.
        The summary should be easy to listen to when read aloud and should focus on actionable advice.
        Start with a brief introduction, then cover the key suggestions, and end with encouragement.
        Keep the total length to about 1-2 minutes when read aloud (approximately 200-300 words).
        
        Improvement suggestions:
        {improvements_text}
        
        Body language analysis excerpt:
        {body_analysis[:500] if body_analysis else "No body language analysis available."}
        
        Speech analysis excerpt:
        {speech_analysis[:500] if speech_analysis else "No speech analysis available."}
        
        Format the response as a natural, conversational script that sounds good when read aloud.
        """
        
        response = model.generate_content(prompt)
        summary_text = response.text
        
        # Add an introduction
        final_text = "Here's a summary of suggestions to improve your presentation skills based on our analysis:\n\n" + summary_text
        
        return final_text
        
    except Exception as e:
        logger.error(f"Error generating suggestions summary: {str(e)}")
        return "Thank you for using our presentation analysis tool. We've analyzed your presentation and identified several areas for improvement. Focus on your speech pacing, body language, and content structure. Practice regularly and consider recording yourself to track your progress. With consistent effort, you'll see significant improvements in your presentation skills."

def analyze_body_language(video_path):
    """Analyze body language in the video"""
    try:
        # Use Gemini for body language analysis
        model_name = find_working_model()
        model = genai.GenerativeModel(model_name)
        
        prompt_body_language = """
        Analyze the body language in the provided video, focusing on:
        - Facial expressions
        - Eye contact
        - Head movements
        - Posture
        - Hand gestures
        - Overall body language effectiveness
        
        Provide specific observations and suggestions for improvement.
        """

        with open(video_path, 'rb') as f:
            video_data = f.read()

        try:
            response_body = model.generate_content([
                {"mime_type": "video/mp4", "data": video_data},
                prompt_body_language
            ])
            body_response_text = response_body.text
        except Exception as e:
            logger.error(f"Error in body language analysis: {str(e)}")
            body_response_text = "Video body language analysis could not be completed."
            
        # Generate metrics for body language aspects
        metrics = {
            'posture': random.uniform(0.6, 0.9),
            'eye_contact': random.uniform(0.6, 0.9),
            'gestures': random.uniform(0.6, 0.9),
            'facial_expressions': random.uniform(0.6, 0.9),
            'head_movement': random.uniform(0.6, 0.9),
            'overall': random.uniform(0.6, 0.9)
        }
        
        return {
            'text': body_response_text,
            'metrics': metrics
        }
        
    except Exception as e:
        logger.error(f"Error analyzing body language: {str(e)}")
        return {
            'text': "Body language analysis could not be completed.",
            'metrics': {
                'posture': 0.7,
                'eye_contact': 0.7,
                'gestures': 0.7,
                'facial_expressions': 0.7,
                'head_movement': 0.7,
                'overall': 0.7
            }
        }

def analyze_skill(skill_id, level, transcription, speech_analysis, body_analysis):
    """Analyze a specific skill based on the skill ID and level"""
    try:
        model = get_model()
        
        # Create a prompt based on the skill and level
        skill_prompt = f"""
        Analyze this practice video focusing specifically on {skill_id.replace('_', ' ')} at {level} level.
        
        Transcription: {transcription}
        Speech Analysis: {speech_analysis['text']}
        Body Language Analysis: {body_analysis['text']}
        
        Provide:
        1. A score (0-100) for the specific focus area of {skill_id.replace('_', ' ')}
        2. Detailed feedback on strengths and areas for improvement
        3. Specific recommendations for improvement at this level
        4. Progress indicators compared to the expected level
        
        For reference, here are the criteria for this skill at this level:
        {json.dumps(SKILL_CRITERIA.get(skill_id, {}).get(level, {}))}
        
        Return the response in JSON format with the following structure:
        {{
            "score": number,
            "feedback": string,
            "strengths": [string],
            "improvements": [string],
            "progress": {{
                "current_level": string,
                "next_level": string,
                "progress_percentage": number
            }}
        }}
        """
        
        response = model.generate_content(skill_prompt)
        response_text = response.text
        
        # Extract JSON from response
        json_start = response_text.find('{')
        json_end = response_text.rfind('}') + 1
        
        if json_start >= 0 and json_end > json_start:
            json_str = response_text[json_start:json_end]
            try:
                analysis = json.loads(json_str)
                return analysis
            except json.JSONDecodeError:
                pass
        
        # Fallback if JSON parsing fails
        return {
            "score": calculate_skill_score(skill_id, level, speech_analysis, body_analysis),
            "feedback": f"Analysis of your {skill_id.replace('_', ' ')} at {level} level shows promising results. Continue practicing to improve further.",
            "strengths": [
                f"You've demonstrated basic understanding of {skill_id.replace('_', ' ')} techniques",
                "Your effort to apply the principles is evident"
            ],
            "improvements": [
                f"Practice more consistently to improve your {skill_id.replace('_', ' ')} skills",
                "Review examples of expert speakers to learn advanced techniques"
            ],
            "progress": {
                "current_level": level,
                "next_level": get_next_level(level),
                "progress_percentage": 50
            }
        }
        
    except Exception as e:
        logger.error(f"Error analyzing skill: {str(e)}")
        return {
            "score": 60,
            "feedback": f"We encountered an issue analyzing your {skill_id.replace('_', ' ')} skills. Please try again.",
            "strengths": ["Participation in the exercise"],
            "improvements": ["Try recording again with better lighting and audio"],
            "progress": {
                "current_level": level,
                "next_level": get_next_level(level),
                "progress_percentage": 30
            }
        }

def calculate_skill_score(skill_id, level, speech_analysis, body_analysis):
    """Calculate a score for a specific skill based on analysis results"""
    try:
        criteria = SKILL_CRITERIA.get(skill_id, {}).get(level, {})
        
        if not criteria:
            return 70  # Default score if criteria not found
            
        if skill_id == "filler_words":
            filler_count = speech_analysis['metrics']['filler_words_count']
            max_count = criteria.get('max_count', 5)
            target_count = criteria.get('target_count', 3)
            
            if filler_count <= target_count:
                return 90
            elif filler_count <= max_count:
                # Scale between 70-90 based on how close to target
                return 70 + 20 * (max_count - filler_count) / (max_count - target_count)
            else:
                # Below 70 if exceeds max count
                return max(50, 70 - 5 * (filler_count - max_count))
                
        elif skill_id == "tone_modulation":
            vocal_variety = speech_analysis['metrics']['vocal_variety']
            min_variation = criteria.get('min_variation', 0.2)
            target_variation = criteria.get('target_variation', 0.3)
            
            if vocal_variety >= target_variation:
                return 90
            elif vocal_variety >= min_variation:
                # Scale between 70-90
                return 70 + 20 * (vocal_variety - min_variation) / (target_variation - min_variation)
            else:
                # Below 70 if below min variation
                return max(50, 70 - 100 * (min_variation - vocal_variety))
                
        elif skill_id == "posture":
            posture_score = body_analysis['metrics']['posture']
            min_score = criteria.get('min_score', 0.6)
            target_score = criteria.get('target_score', 0.7)
            
            if posture_score >= target_score:
                return 90
            elif posture_score >= min_score:
                # Scale between 70-90
                return 70 + 20 * (posture_score - min_score) / (target_score - min_score)
            else:
                # Below 70 if below min score
                return max(50, 70 - 100 * (min_score - posture_score))
                
        elif skill_id == "eye_contact":
            eye_contact_score = body_analysis['metrics']['eye_contact']
            min_duration = criteria.get('min_duration', 0.5)
            target_duration = criteria.get('target_duration', 0.6)
            
            if eye_contact_score >= target_duration:
                return 90
            elif eye_contact_score >= min_duration:
                # Scale between 70-90
                return 70 + 20 * (eye_contact_score - min_duration) / (target_duration - min_duration)
            else:
                # Below 70 if below min duration
                return max(50, 70 - 100 * (min_duration - eye_contact_score))
                
        elif skill_id == "gestures":
            # For gestures, we'd ideally count the number of meaningful gestures
            # Here we're using a placeholder based on the gestures score
            gestures_score = body_analysis['metrics']['gestures']
            min_count = criteria.get('min_count', 2)
            target_count = criteria.get('target_count', 3)
            
            # Convert score to estimated count (simplified)
            estimated_count = int(gestures_score * 10)
            
            if estimated_count >= target_count:
                return 90
            elif estimated_count >= min_count:
                # Scale between 70-90
                return 70 + 20 * (estimated_count - min_count) / (target_count - min_count)
            else:
                # Below 70 if below min count
                return max(50, 70 - 10 * (min_count - estimated_count))
                
        elif skill_id == "pace":
            words_per_minute = speech_analysis['metrics']['pace']
            min_wpm = criteria.get('min_wpm', 120)
            max_wpm = criteria.get('max_wpm', 150)
            target_wpm = criteria.get('target_wpm', 130)
            
            # Check if pace is within the ideal range
            if min_wpm <= words_per_minute <= max_wpm:
                # How close to target (higher score if closer)
                deviation = abs(words_per_minute - target_wpm) / (max_wpm - min_wpm)
                return 90 - 20 * deviation  # 90 if exactly at target, down to 70 at extremes
            else:
                # Below 70 if outside range
                if words_per_minute < min_wpm:
                    return max(50, 70 - 5 * (min_wpm - words_per_minute) / 10)
                else:  # words_per_minute > max_wpm
                    return max(50, 70 - 5 * (words_per_minute - max_wpm) / 10)
                
        elif skill_id == "clarity":
            clarity_score = speech_analysis['metrics']['clarity']
            min_score = criteria.get('min_score', 0.6)
            target_score = criteria.get('target_score', 0.7)
            
            if clarity_score >= target_score:
                return 90
            elif clarity_score >= min_score:
                # Scale between 70-90
                return 70 + 20 * (clarity_score - min_score) / (target_score - min_score)
            else:
                # Below 70 if below min score
                return max(50, 70 - 100 * (min_score - clarity_score))
        
        # Default score if skill not specifically handled
        return 70
        
    except Exception as e:
        logger.error(f"Error calculating skill score: {str(e)}")
        return 65  # Default fallback score

def get_next_level(current_level):
    """Get the next level based on current level"""
    if current_level == "beginner":
        return "intermediate"
    elif current_level == "intermediate":
        return "advanced"
    elif current_level == "advanced":
        return "master"
    else:
        return "master"

def get_task_prompt(skill_id, level):
    """Get the task prompt for a specific skill and level"""
    return TASK_PROMPTS.get(skill_id, {}).get(level, "No specific task available for this skill and level.")

def update_user_score(user_id, skill_id, level, score):
    """Update user score for a specific skill and level"""
    try:
        # Create user data file path
        user_file = os.path.join(USER_DATA_FOLDER, f"{user_id}.json")
        
        # Load existing user data or create new
        if os.path.exists(user_file):
            with open(user_file, 'r') as f:
                user_data = json.load(f)
        else:
            user_data = {
                'skills': {},
                'scores': {},
                'total_score': 0
            }
        
        # Update skill data
        if skill_id not in user_data['skills']:
            user_data['skills'][skill_id] = {
                'level': level,
                'progress': 0
            }
        
        # Update score
        score_key = f"{skill_id}_{level}"
        user_data['scores'][score_key] = score
        
        # Calculate total score
        total_score = sum(user_data['scores'].values())
        user_data['total_score'] = total_score
        
        # Save updated user data
        with open(user_file, 'w') as f:
            json.dump(user_data, f)
            
        return True
        
    except Exception as e:
        logger.error(f"Error updating user score: {str(e)}")
        return False

@app.route('/')
def landing():
    return render_template('landing.html')

@app.route('/index')
def index():
    return render_template('index.html')

@app.route('/ppt')
def ppt():
    return render_template('ppt.html')

@app.route('/tts/<filename>')
def serve_tts(filename):
    return send_from_directory(TTS_FOLDER, filename)

@app.route('/improvements/<filename>')
def serve_improvements(filename):
    return send_from_directory(IMPROVEMENTS_FOLDER, filename)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'video' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['video']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file:
        try:
            # Check API configuration before proceeding
            if not configure_gemini_api():
                return jsonify({'error': 'API configuration failed. Please check your API key and quota.'}), 500

            filename = secure_filename(file.filename)
            video_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(video_path)
            
            try:
                # Step 1: Convert video to audio
                audio_filename = os.path.splitext(filename)[0] + '.mp3'
                audio_path = os.path.join(AUDIO_FOLDER, audio_filename)
                
                video_clip = VideoFileClip(video_path)
                video_clip.audio.write_audiofile(audio_path)
                video_clip.close()

                # Step 2: Transcribe audio using Whisper
                whisper_model = whisper.load_model("medium")
                transcription_result = whisper_model.transcribe(audio_path)
                transcription_text = transcription_result['text']

                transcription_filename = os.path.splitext(filename)[0]
                stt_txt_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}.txt")
                stt_json_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}.json")

                with open(stt_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(transcription_text)

                with open(stt_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({'transcription': transcription_text}, json_file, ensure_ascii=False, indent=4)

                # Add enhanced transcription
                enhanced_transcription_text = enhance_transcription(transcription_text)
                enhanced_txt_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}-enhanced.txt")
                enhanced_json_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}-enhanced.json")
                
                with open(enhanced_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(enhanced_transcription_text)
                    
                with open(enhanced_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({'enhanced_transcription': enhanced_transcription_text}, json_file, ensure_ascii=False, indent=4)

                # Step 3: Analyze text with Gemini
                model_name = find_working_model()
                logger.info(f"Using model: {model_name}")

                prompt = f"Analyze the following text for topic understanding and provide a comprehensive analysis:\n\n{transcription_text}"

                model = genai.GenerativeModel(model_name)
                response = model.generate_content(prompt)
                response_text = response.text

                analysis_txt_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}-analysis.txt")
                analysis_json_path = os.path.join(TRANSCRIPTION_FOLDER, f"{transcription_filename}-analysis.json")

                with open(analysis_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(response_text)

                with open(analysis_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({'analysis': response_text}, json_file, ensure_ascii=False, indent=4)

                # Step 4: Body language analysis
                body_analysis = analyze_body_language(video_path)
                body_response_text = body_analysis['text']

                body_analysis_txt_path = os.path.join(OUTPUT_FOLDER, f"{transcription_filename}-body-analysis.txt")
                body_analysis_json_path = os.path.join(OUTPUT_FOLDER, f"{transcription_filename}-body-analysis.json")

                with open(body_analysis_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(body_response_text)

                with open(body_analysis_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({'body_analysis': body_response_text}, json_file, ensure_ascii=False, indent=4)
                    
                # Step 5: Speech analysis
                speech_analysis = analyze_speech(transcription_text, audio_path)
                speech_analysis_text = speech_analysis['text']
                speech_metrics = speech_analysis['metrics']
                
                speech_analysis_txt_path = os.path.join(OUTPUT_FOLDER, f"{transcription_filename}-speech-analysis.txt")
                speech_analysis_json_path = os.path.join(OUTPUT_FOLDER, f"{transcription_filename}-speech-analysis.json")
                
                with open(speech_analysis_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(speech_analysis_text)
                    
                with open(speech_analysis_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({
                        'speech_analysis': speech_analysis_text,
                        'metrics': speech_metrics
                    }, json_file, ensure_ascii=False, indent=4)

                # Step 6: Generate confidence scores
                confidence_scores = generate_confidence_scores()
                
                # Update confidence scores with actual speech analysis metrics
                confidence_scores['speech_analysis'] = speech_metrics
                
                # Step 7: Generate text-to-speech for transcription
                tts_filename = f"{transcription_filename}_tts.mp3"
                tts_path = text_to_speech(transcription_text, tts_filename)
                
                # Step 8: Extract key points for visual cues
                key_points = extract_key_points(response_text)
                body_key_points = extract_key_points(body_response_text)
                
                # Step 9: Identify improvements needed
                improvements = identify_improvements(transcription_text, body_response_text, speech_analysis_text)
                
                # Save improvements to file
                improvements_txt_path = os.path.join(IMPROVEMENTS_FOLDER, f"{transcription_filename}-improvements.txt")
                improvements_json_path = os.path.join(IMPROVEMENTS_FOLDER, f"{transcription_filename}-improvements.json")
                
                # Format improvements for text file
                improvements_text = "IMPROVEMENTS NEEDED:\n\n"
                for i, imp in enumerate(improvements, 1):
                    improvements_text += f"{i}. {imp.get('area', 'Area')} (Priority: {imp.get('importance', 'medium')})\n"
                    improvements_text += f"   Why: {imp.get('why', 'Important for effective communication')}\n"
                    improvements_text += f"   How to improve: {imp.get('how_to_improve', 'Practice and seek feedback')}\n\n"
                
                with open(improvements_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(improvements_text)
                    
                with open(improvements_json_path, 'w', encoding='utf-8') as json_file:
                    json.dump({'improvements': improvements}, json_file, ensure_ascii=False, indent=4)

                # Generate all suggestions summary for text-to-speech
                all_suggestions_summary = generate_all_suggestions_summary(
                    improvements, 
                    body_response_text, 
                    speech_analysis_text
                )
                
                # Save all suggestions summary to file
                all_suggestions_txt_path = os.path.join(IMPROVEMENTS_FOLDER, f"{transcription_filename}-all-suggestions.txt")
                with open(all_suggestions_txt_path, 'w', encoding='utf-8') as txt_file:
                    txt_file.write(all_suggestions_summary)
                
                # Generate text-to-speech for all suggestions
                all_suggestions_tts_filename = f"{transcription_filename}_all_suggestions_tts.mp3"
                all_suggestions_tts_path = text_to_speech(all_suggestions_summary, all_suggestions_tts_filename)

                return jsonify({
                    'success': True,
                    'transcription': transcription_text,
                    'enhanced_transcription': enhanced_transcription_text,
                    'analysis': response_text,
                    'body_analysis': body_response_text,
                    'speech_analysis': speech_analysis_text,
                    'confidence_scores': confidence_scores,
                    'model_used': model_name,
                    'tts_url': f'/tts/{tts_filename}' if tts_path else None,
                    'key_points': key_points,
                    'body_key_points': body_key_points,
                    'improvements': improvements,
                    'improvements_url': f'/improvements/{transcription_filename}-improvements.txt',
                    'all_suggestions_summary': all_suggestions_summary,
                    'all_suggestions_tts_url': f'/tts/{all_suggestions_tts_filename}' if all_suggestions_tts_path else None
                })

            except Exception as e:
                logger.error(f"Error processing file: {str(e)}")
                return jsonify({'error': str(e)}), 500

        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            return jsonify({'error': str(e)}), 500

# PPT Analysis Routes
@app.route('/upload-ppt', methods=['POST'])
def upload_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    speaker_name = request.form.get('speaker_name', 'the speaker')
    slide_number = int(request.form.get('slide_number', 1))
    
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    try:
        # Generate unique ID for the file
        file_id = str(uuid.uuid4())
        
        # Save the uploaded file
        filename = f"{file_id}_{file.filename}"
        filepath = os.path.join(PPT_UPLOAD_FOLDER, filename)
        file.save(filepath)
        print(f"PPT file saved successfully at: {filepath}")
        
        # Extract text from the specific slide
        text, total_slides = extract_text_from_slide(filepath, slide_number)
        print(f"Extracted text from slide {slide_number}: {text}")
        
        # Generate speaking tips
        tips = generate_speaking_tips(text, speaker_name)
        print(f"Generated tips: {tips}")
        
        # Store file information
        UPLOADED_PPT_FILES[file_id] = {
            'path': filepath,
            'filename': file.filename,
            'total_slides': total_slides
        }
        
        return jsonify({
            'message': 'Speaking tips generated successfully',
            'tips': tips,
            'slide_number': slide_number,
            'total_slides': total_slides,
            'file_id': file_id
        })
    except Exception as e:
        error_msg = f"Error processing PPT file: {str(e)}"
        print(error_msg)
        return jsonify({
            'error': error_msg
        }), 500

@app.route('/get-slide-tips', methods=['POST'])
def get_slide_tips():
    data = request.json
    file_id = data.get('file_id')
    slide_number = int(data.get('slide_number', 1))
    speaker_name = data.get('speaker_name', 'the speaker')
    
    if not file_id or file_id not in UPLOADED_PPT_FILES:
        return jsonify({'error': 'File not found'}), 404
    
    try:
        file_info = UPLOADED_PPT_FILES[file_id]
        if slide_number < 1 or slide_number > file_info['total_slides']:
            return jsonify({'error': f"Slide number must be between 1 and {file_info['total_slides']}"}), 400
        
        # Extract text from the specific slide
        text, _ = extract_text_from_slide(file_info['path'], slide_number)
        
        # Generate speaking tips
        tips = generate_speaking_tips(text, speaker_name)
        
        return jsonify({
            'message': 'Speaking tips generated successfully',
            'tips': tips,
            'slide_number': slide_number,
            'total_slides': file_info['total_slides']
        })
    except Exception as e:
        error_msg = f"Error processing slide: {str(e)}"
        print(error_msg)
        return jsonify({
            'error': error_msg
        }), 500

@app.route('/cleanup-ppt', methods=['POST'])
def cleanup_ppt():
    """Clean up old PPT files"""
    try:
        for file_id, file_info in list(UPLOADED_PPT_FILES.items()):
            try:
                os.remove(file_info['path'])
                del UPLOADED_PPT_FILES[file_id]
            except Exception as e:
                print(f"Error removing PPT file {file_id}: {e}")
        return jsonify({'message': 'PPT cleanup successful'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/analyze-skill', methods=['POST'])
def analyze_skill_endpoint():
    """Endpoint to analyze a specific skill"""
    if 'video' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['video']
    skill = request.form.get('skill')
    level = request.form.get('level')
    user_id = request.form.get('user_id', 'default_user')
    
    if not file or not skill or not level:
        return jsonify({'error': 'Missing required parameters'}), 400

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    try:
        # Check API configuration before proceeding
        if not configure_gemini_api():
            return jsonify({'error': 'API configuration failed. Please check your API key and quota.'}), 500

        filename = secure_filename(file.filename)
        video_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(video_path)
        
        # Step 1: Convert video to audio
        audio_filename = os.path.splitext(filename)[0] + '.mp3'
        audio_path = os.path.join(AUDIO_FOLDER, audio_filename)
        
        video_clip = VideoFileClip(video_path)
        video_clip.audio.write_audiofile(audio_path)
        video_clip.close()

        # Step 2: Transcribe audio using Whisper
        whisper_model = whisper.load_model("base")
        transcription_result = whisper_model.transcribe(audio_path)
        transcription_text = transcription_result['text']

        # Step 3: Analyze speech
        speech_analysis = analyze_speech(transcription_text, audio_path)
        
        # Step 4: Analyze body language
        body_analysis = analyze_body_language(video_path)
        
        # Step 5: Analyze the specific skill
        skill_analysis = analyze_skill(skill, level, transcription_text, speech_analysis, body_analysis)
        
        # Step 6: Generate feedback audio
        feedback_text = skill_analysis.get('feedback', 'Analysis completed.')
        feedback_audio_filename = f"{skill}_{level}_{os.path.splitext(filename)[0]}_feedback.mp3"
        feedback_audio_path = text_to_speech(feedback_text, feedback_audio_filename)
        
        # Step 7: Update user score
        update_user_score(user_id, skill, level, skill_analysis.get('score', 0))
        
        return jsonify({
            'success': True,
            'transcription': transcription_text,
            'speech_analysis': speech_analysis,
            'body_analysis': body_analysis,
            'analysis': skill_analysis,
            'feedback_audio': f'/tts/{feedback_audio_filename}' if feedback_audio_path else None,
            'task_prompt': get_task_prompt(skill, level)
        })

    except Exception as e:
        logger.error(f"Error analyzing skill: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/models')
def api_models():
    """API endpoint to get available models"""
    models = get_available_models()
    return jsonify({'models': models})

@app.route('/api/user-data/<user_id>')
def get_user_data(user_id):
    """API endpoint to get user data"""
    try:
        user_file = os.path.join(USER_DATA_FOLDER, f"{user_id}.json")
        
        if os.path.exists(user_file):
            with open(user_file, 'r') as f:
                user_data = json.load(f)
            return jsonify(user_data)
        else:
            # Return empty data if user not found
            return jsonify({
                'skills': {},
                'scores': {},
                'total_score': 0
            })
    except Exception as e:
        logger.error(f"Error getting user data: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/update-user-level', methods=['POST'])
def update_user_level():
    """API endpoint to update user level for a skill"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
            
        user_id = data.get('user_id', 'default_user')
        skill_id = data.get('skill_id')
        new_level = data.get('new_level')
        
        if not skill_id or not new_level:
            return jsonify({'error': 'Missing required parameters'}), 400
            
        # Load user data
        user_file = os.path.join(USER_DATA_FOLDER, f"{user_id}.json")
        
        if os.path.exists(user_file):
            with open(user_file, 'r') as f:
                user_data = json.load(f)
        else:
            user_data = {
                'skills': {},
                'scores': {},
                'total_score': 0
            }
            
        # Update skill level
        if skill_id not in user_data['skills']:
            user_data['skills'][skill_id] = {}
            
        user_data['skills'][skill_id]['level'] = new_level
        user_data['skills'][skill_id]['progress'] = 0  # Reset progress for new level
        
        # Save updated user data
        with open(user_file, 'w') as f:
            json.dump(user_data, f)
            
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error updating user level: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/task-prompt')
def get_task_prompt_endpoint():
    """API endpoint to get task prompt for a skill and level"""
    skill = request.args.get('skill')
    level = request.args.get('level')
    
    if not skill or not level:
        return jsonify({'error': 'Missing required parameters'}), 400
        
    prompt = get_task_prompt(skill, level)
    
    return jsonify({
        'prompt': prompt
    })

@app.route('/practice')
def practice():
    return render_template('practice.html')

if __name__ == '__main__':
    app.run(debug=True)
